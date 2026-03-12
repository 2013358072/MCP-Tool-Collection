from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional

import requests
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
from fastmcp import FastMCP

from mcp_toolkit.core import config as _cfg
from mcp_toolkit.providers.base import BaseProvider

# ======================================================================
# 本模块仅处理 .pptx 格式演示文稿（python-pptx）
# ======================================================================

_CHART_TYPE_MAP: Dict[str, Any] = {
    "bar":             XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked":     XL_CHART_TYPE.BAR_STACKED,
    "column":          XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked":  XL_CHART_TYPE.COLUMN_STACKED,
    "line":            XL_CHART_TYPE.LINE,
    "line_marked":     XL_CHART_TYPE.LINE_MARKERS,
    "pie":             XL_CHART_TYPE.PIE,
    "area":            XL_CHART_TYPE.AREA,
    "scatter":         XL_CHART_TYPE.XY_SCATTER,
}

_ALIGN_MAP: Dict[str, Any] = {
    "left":    PP_ALIGN.LEFT,
    "center":  PP_ALIGN.CENTER,
    "right":   PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

# 16:9 宽屏默认尺寸（厘米）
_DEFAULT_WIDTH_CM  = 33.87
_DEFAULT_HEIGHT_CM = 19.05


# ======================================================================
# 内部工具函数
# ======================================================================

def _mkparent(p: Path) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)


def _sandbox_root() -> Path:
    return Path(_cfg.FILESYSTEM_ROOT or Path.cwd()).resolve()


def _resolve_slide_path(path: str) -> Path:
    """将 PPT 路径限制在沙箱目录内。"""
    sandbox_root = _sandbox_root()
    p = Path(path)
    resolved = (sandbox_root / p).resolve() if not p.is_absolute() else p.resolve()
    try:
        resolved.relative_to(sandbox_root)
    except ValueError as e:
        raise PermissionError(
            f"路径越出沙箱: {resolved}（FILESYSTEM_ROOT={sandbox_root}）"
        ) from e
    return resolved


def _check_pptx(p: Path) -> Optional[Dict[str, Any]]:
    if not p.exists() or not p.is_file():
        return {"ok": False, "error": "NOT_FOUND", "path": str(p)}
    if p.suffix.lower() != ".pptx":
        return {"ok": False, "error": "NOT_PPTX", "detail": f"仅支持 .pptx 格式，收到: {p.suffix}"}
    return None


def _rgb(color: str) -> RGBColor:
    c = color.lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _emu_to_cm(emu: int) -> float:
    return round(emu / 914400 * 2.54, 2)


def _fetch_image_local(image_path: str) -> tuple[str, bool]:
    """返回 (本地路径, 是否为临时文件)。URL 自动下载为临时文件。"""
    if image_path.startswith(("http://", "https://")):
        resp = requests.get(image_path, timeout=30)
        resp.raise_for_status()
        suffix = Path(image_path.split("?")[0]).suffix or ".png"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(resp.content)
        tmp.close()
        return tmp.name, True
    return str(_resolve_slide_path(image_path)), False


# ---------------------------------------------------------------------- #
# slides_create_deck
# ---------------------------------------------------------------------- #

def _slides_create_deck(
    path: str,
    title: Optional[str] = None,
    subtitle: Optional[str] = None,
    width_cm: float = _DEFAULT_WIDTH_CM,
    height_cm: float = _DEFAULT_HEIGHT_CM,
) -> Dict[str, Any]:
    """创建新的 .pptx 文件，可选带标题首页。"""
    try:
        p = _resolve_slide_path(path)
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if p.suffix.lower() != ".pptx":
        return {"ok": False, "error": "NOT_PPTX", "detail": f"路径必须以 .pptx 结尾，收到: {p.suffix}"}
    try:
        _mkparent(p)
        prs = Presentation()
        prs.slide_width  = Cm(width_cm)
        prs.slide_height = Cm(height_cm)

        if title:
            layout = prs.slide_layouts[0]   # 标题幻灯片版式
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = title
            if subtitle:
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph.text = subtitle
                        break

        prs.save(p)
        return {
            "ok": True,
            "path": str(p),
            "slide_count": len(prs.slides),
            "width_cm": width_cm,
            "height_cm": height_cm,
        }
    except Exception as e:
        return {"ok": False, "error": "CREATE_FAILED", "path": str(p), "detail": str(e)}


# ---------------------------------------------------------------------- #
# slides_add_slide
# ---------------------------------------------------------------------- #

def _slides_add_slide(
    path: str,
    layout_index: int = 6,
    title: Optional[str] = None,
    subtitle: Optional[str] = None,
) -> Dict[str, Any]:
    """新增一页幻灯片。layout_index 参考内置版式序号，6 = 空白。"""
    try:
        p = _resolve_slide_path(path)
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if err := _check_pptx(p):
        return err
    try:
        prs = Presentation(p)
        idx = max(0, min(layout_index, len(prs.slide_layouts) - 1))
        layout = prs.slide_layouts[idx]
        slide = prs.slides.add_slide(layout)

        if title and slide.shapes.title:
            slide.shapes.title.text = title
        if subtitle:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = subtitle
                    break

        prs.save(p)
        slide_index = len(prs.slides) - 1
        return {
            "ok": True,
            "path": str(p),
            "slide_index": slide_index,
            "slide_count": len(prs.slides),
            "layout_index": idx,
        }
    except Exception as e:
        return {"ok": False, "error": "ADD_SLIDE_FAILED", "path": str(p), "detail": str(e)}


# ---------------------------------------------------------------------- #
# slides_add_text
# ---------------------------------------------------------------------- #

def _slides_add_text(
    path: str,
    slide_index: int,
    text: str,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    font_size_pt: float = 18.0,
    font_name: str = "Arial",
    font_bold: bool = False,
    font_italic: bool = False,
    font_underline: bool = False,
    font_color: Optional[str] = None,
    alignment: str = "left",
    background_color: Optional[str] = None,
    word_wrap: bool = True,
) -> Dict[str, Any]:
    """在指定幻灯片添加文本框。text 中的换行符 \\n 会被拆分为多个段落。"""
    try:
        p = _resolve_slide_path(path)
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if err := _check_pptx(p):
        return err
    try:
        prs = Presentation(p)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return {"ok": False, "error": "SLIDE_NOT_FOUND", "slide_index": slide_index}

        slide = prs.slides[slide_index]
        txBox = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))

        tf = txBox.text_frame
        tf.word_wrap = word_wrap

        if background_color:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = _rgb(background_color)

        align = _ALIGN_MAP.get(alignment.lower(), PP_ALIGN.LEFT)
        lines = text.split("\n")

        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            run = para.add_run()
            run.text = line
            run.font.size      = Pt(font_size_pt)
            run.font.name      = font_name
            run.font.bold      = font_bold
            run.font.italic    = font_italic
            run.font.underline = font_underline
            if font_color:
                run.font.color.rgb = _rgb(font_color)

        prs.save(p)
        return {
            "ok": True,
            "path": str(p),
            "slide_index": slide_index,
            "shape_name": txBox.name,
            "position": {"left_cm": left_cm, "top_cm": top_cm},
            "size": {"width_cm": width_cm, "height_cm": height_cm},
        }
    except Exception as e:
        return {"ok": False, "error": "ADD_TEXT_FAILED", "path": str(p), "detail": str(e)}


# ---------------------------------------------------------------------- #
# slides_add_image
# ---------------------------------------------------------------------- #

def _slides_add_image(
    path: str,
    slide_index: int,
    image_path: str,
    left_cm: float,
    top_cm: float,
    width_cm: Optional[float] = None,
    height_cm: Optional[float] = None,
) -> Dict[str, Any]:
    """在指定幻灯片插入图片（本地路径或 URL）。
    width_cm 和 height_cm 均为空时保持原始比例。"""
    try:
        p = _resolve_slide_path(path)
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if err := _check_pptx(p):
        return err
    tmp_path, is_tmp = None, False
    try:
        prs = Presentation(p)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return {"ok": False, "error": "SLIDE_NOT_FOUND", "slide_index": slide_index}

        img_local, is_tmp = _fetch_image_local(image_path)
        tmp_path = img_local if is_tmp else None

        slide = prs.slides[slide_index]
        kwargs: Dict[str, Any] = {
            "image_file": img_local,
            "left": Cm(left_cm),
            "top":  Cm(top_cm),
        }
        if width_cm is not None:
            kwargs["width"] = Cm(width_cm)
        if height_cm is not None:
            kwargs["height"] = Cm(height_cm)

        pic = slide.shapes.add_picture(**kwargs)
        prs.save(p)
        return {
            "ok": True,
            "path": str(p),
            "slide_index": slide_index,
            "shape_name": pic.name,
            "position": {"left_cm": left_cm, "top_cm": top_cm},
            "size": {
                "width_cm":  _emu_to_cm(pic.width),
                "height_cm": _emu_to_cm(pic.height),
            },
        }
    except Exception as e:
        return {"ok": False, "error": "ADD_IMAGE_FAILED", "path": str(p), "detail": str(e)}
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)


# ---------------------------------------------------------------------- #
# slides_add_table
# ---------------------------------------------------------------------- #

def _slides_add_table(
    path: str,
    slide_index: int,
    rows: List[List[str]],
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    has_header: bool = True,
    font_size_pt: float = 12.0,
    font_name: str = "Arial",
    header_font_bold: bool = True,
    header_fill_color: Optional[str] = None,
    header_font_color: Optional[str] = None,
    col_widths_cm: Optional[List[float]] = None,
) -> Dict[str, Any]:
    """在指定幻灯片插入表格。rows 为二维字符串数组，has_header=True 时首行加粗。"""
    try:
        p = _resolve_slide_path(path)
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if err := _check_pptx(p):
        return err
    if not rows:
        return {"ok": False, "error": "EMPTY_ROWS", "detail": "rows 不能为空"}
    try:
        prs = Presentation(p)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return {"ok": False, "error": "SLIDE_NOT_FOUND", "slide_index": slide_index}

        slide = prs.slides[slide_index]
        row_count = len(rows)
        col_count = max(len(r) for r in rows)

        tbl_shape = slide.shapes.add_table(
            row_count, col_count,
            Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm),
        )
        tbl = tbl_shape.table

        if col_widths_cm and len(col_widths_cm) == col_count:
            for i, cw in enumerate(col_widths_cm):
                tbl.columns[i].width = Cm(cw)

        for r_idx, row in enumerate(rows):
            is_header = has_header and r_idx == 0
            for c_idx in range(col_count):
                cell_text = str(row[c_idx]) if c_idx < len(row) else ""
                cell = tbl.cell(r_idx, c_idx)
                tf = cell.text_frame

                # 清空并重新写入，确保样式可控
                tf.clear()
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = cell_text
                run.font.size = Pt(font_size_pt)
                run.font.name = font_name

                if is_header:
                    run.font.bold = header_font_bold
                    if header_font_color:
                        run.font.color.rgb = _rgb(header_font_color)
                    if header_fill_color:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(header_fill_color)

        prs.save(p)
        return {
            "ok": True,
            "path": str(p),
            "slide_index": slide_index,
            "shape_name": tbl_shape.name,
            "rows": row_count,
            "cols": col_count,
            "position": {"left_cm": left_cm, "top_cm": top_cm},
            "size": {"width_cm": width_cm, "height_cm": height_cm},
        }
    except Exception as e:
        return {"ok": False, "error": "ADD_TABLE_FAILED", "path": str(p), "detail": str(e)}


# ---------------------------------------------------------------------- #
# slides_add_chart
# ---------------------------------------------------------------------- #

def _slides_add_chart(
    path: str,
    slide_index: int,
    chart_type: str,
    categories: List[str],
    series: List[Dict[str, Any]],
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    title: Optional[str] = None,
    has_legend: bool = True,
    has_data_labels: bool = False,
) -> Dict[str, Any]:
    """在指定幻灯片插入图表。

    chart_type 可选: bar / bar_stacked / column / column_stacked /
                     line / line_marked / pie / area / scatter
    series 每项格式: {"name": "系列名", "values": [1.0, 2.0, 3.0]}
    """
    try:
        p = _resolve_slide_path(path)
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if err := _check_pptx(p):
        return err
    chart_enum = _CHART_TYPE_MAP.get(chart_type.lower())
    if chart_enum is None:
        return {
            "ok": False,
            "error": "UNKNOWN_CHART_TYPE",
            "supported": list(_CHART_TYPE_MAP.keys()),
        }
    try:
        prs = Presentation(p)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return {"ok": False, "error": "SLIDE_NOT_FOUND", "slide_index": slide_index}

        slide = prs.slides[slide_index]
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s in series:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        chart_frame = slide.shapes.add_chart(
            chart_enum,
            Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm),
            chart_data,
        )
        chart = chart_frame.chart

        if title:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = title

        chart.has_legend = has_legend

        if has_data_labels and chart.plots:
            chart.plots[0].has_data_labels = True

        prs.save(p)
        return {
            "ok": True,
            "path": str(p),
            "slide_index": slide_index,
            "shape_name": chart_frame.name,
            "chart_type": chart_type,
            "series_count": len(series),
            "category_count": len(categories),
            "position": {"left_cm": left_cm, "top_cm": top_cm},
            "size": {"width_cm": width_cm, "height_cm": height_cm},
        }
    except Exception as e:
        return {"ok": False, "error": "ADD_CHART_FAILED", "path": str(p), "detail": str(e)}


# ---------------------------------------------------------------------- #
# slides_export_pptx
# ---------------------------------------------------------------------- #

def _slides_export_pptx(
    path: str,
    output_path: Optional[str] = None,
) -> Dict[str, Any]:
    """将演示文稿导出/另存为 pptx；output_path 为空时原路径覆盖保存。"""
    try:
        p = _resolve_slide_path(path)
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if err := _check_pptx(p):
        return err
    try:
        out = _resolve_slide_path(output_path) if output_path else p
    except PermissionError as e:
        return {"ok": False, "error": "PATH_OUTSIDE_SANDBOX", "detail": str(e)}
    if out.suffix.lower() != ".pptx":
        return {"ok": False, "error": "NOT_PPTX", "detail": f"输出路径必须以 .pptx 结尾，收到: {out.suffix}"}
    try:
        prs = Presentation(p)
        _mkparent(out)
        prs.save(out)
        return {
            "ok": True,
            "path": str(out),
            "slide_count": len(prs.slides),
            "size_bytes": out.stat().st_size,
        }
    except Exception as e:
        return {"ok": False, "error": "EXPORT_FAILED", "path": str(p), "detail": str(e)}


# ======================================================================
# Provider
# ======================================================================

class SlidesProvider(BaseProvider):
    """PowerPoint 演示文稿工具集 Provider（仅处理 .pptx 格式，依赖 python-pptx）。"""

    @property
    def name(self) -> str:
        return "slides"

    def is_available(self) -> bool:
        return True

    async def initialize(self) -> None:
        self.logger.local("info", "SlidesProvider 初始化完成")

    def register(self, mcp: FastMCP) -> None:

        @mcp.tool()
        async def slides_create_deck(
            path: str,
            title: Optional[str] = None,
            subtitle: Optional[str] = None,
            width_cm: float = _DEFAULT_WIDTH_CM,
            height_cm: float = _DEFAULT_HEIGHT_CM,
        ) -> Dict[str, Any]:
            """创建新的 .pptx 演示文稿。
            width_cm / height_cm: 幻灯片尺寸（厘米），默认 16:9 宽屏（33.87×19.05）
            title / subtitle: 若提供则自动添加标题首页
            """
            return _slides_create_deck(path, title=title, subtitle=subtitle,
                                       width_cm=width_cm, height_cm=height_cm)

        @mcp.tool()
        async def slides_add_slide(
            path: str,
            layout_index: int = 6,
            title: Optional[str] = None,
            subtitle: Optional[str] = None,
        ) -> Dict[str, Any]:
            """向已有 pptx 新增一页幻灯片，返回新页的 slide_index（0-based）。
            layout_index: 内置版式序号（0=标题页, 1=标题+内容, 6=空白）
            """
            return _slides_add_slide(path, layout_index=layout_index,
                                     title=title, subtitle=subtitle)

        @mcp.tool()
        async def slides_add_text(
            path: str,
            slide_index: int,
            text: str,
            left_cm: float,
            top_cm: float,
            width_cm: float,
            height_cm: float,
            font_size_pt: float = 18.0,
            font_name: str = "Arial",
            font_bold: bool = False,
            font_italic: bool = False,
            font_underline: bool = False,
            font_color: Optional[str] = None,
            alignment: str = "left",
            background_color: Optional[str] = None,
            word_wrap: bool = True,
        ) -> Dict[str, Any]:
            """在指定幻灯片添加文本框（所有位置/尺寸单位为厘米）。
            font_color / background_color: RRGGBB 16进制颜色，如 "FF0000"
            alignment: left / center / right / justify
            text 中的 \\n 会被拆分为多个段落
            """
            return _slides_add_text(
                path, slide_index, text,
                left_cm=left_cm, top_cm=top_cm, width_cm=width_cm, height_cm=height_cm,
                font_size_pt=font_size_pt, font_name=font_name,
                font_bold=font_bold, font_italic=font_italic, font_underline=font_underline,
                font_color=font_color, alignment=alignment,
                background_color=background_color, word_wrap=word_wrap,
            )

        @mcp.tool()
        async def slides_add_image(
            path: str,
            slide_index: int,
            image_path: str,
            left_cm: float,
            top_cm: float,
            width_cm: Optional[float] = None,
            height_cm: Optional[float] = None,
        ) -> Dict[str, Any]:
            """在指定幻灯片插入图片（本地路径或 HTTP/HTTPS URL）。
            width_cm / height_cm: 均为空时保持图片原始比例；仅指定一个时等比缩放
            """
            return _slides_add_image(path, slide_index, image_path,
                                     left_cm=left_cm, top_cm=top_cm,
                                     width_cm=width_cm, height_cm=height_cm)

        @mcp.tool()
        async def slides_add_table(
            path: str,
            slide_index: int,
            rows: List[List[str]],
            left_cm: float,
            top_cm: float,
            width_cm: float,
            height_cm: float,
            has_header: bool = True,
            font_size_pt: float = 12.0,
            font_name: str = "Arial",
            header_font_bold: bool = True,
            header_fill_color: Optional[str] = None,
            header_font_color: Optional[str] = None,
            col_widths_cm: Optional[List[float]] = None,
        ) -> Dict[str, Any]:
            """在指定幻灯片插入表格。rows 为二维字符串数组。
            has_header: 首行作为表头（加粗、可设背景色）
            col_widths_cm: 各列宽度（厘米），长度须与列数相等
            header_fill_color / header_font_color: RRGGBB 16进制颜色
            """
            return _slides_add_table(
                path, slide_index, rows,
                left_cm=left_cm, top_cm=top_cm, width_cm=width_cm, height_cm=height_cm,
                has_header=has_header, font_size_pt=font_size_pt, font_name=font_name,
                header_font_bold=header_font_bold,
                header_fill_color=header_fill_color, header_font_color=header_font_color,
                col_widths_cm=col_widths_cm,
            )

        @mcp.tool()
        async def slides_add_chart(
            path: str,
            slide_index: int,
            chart_type: str,
            categories: List[str],
            series: List[Dict[str, Any]],
            left_cm: float,
            top_cm: float,
            width_cm: float,
            height_cm: float,
            title: Optional[str] = None,
            has_legend: bool = True,
            has_data_labels: bool = False,
        ) -> Dict[str, Any]:
            """在指定幻灯片插入图表。
            chart_type: bar / bar_stacked / column / column_stacked /
                        line / line_marked / pie / area / scatter
            categories: X 轴类别列表，如 ["Q1", "Q2", "Q3"]
            series: 数据系列列表，每项 {"name": "销售额", "values": [100, 200, 150]}
            """
            return _slides_add_chart(
                path, slide_index, chart_type, categories, series,
                left_cm=left_cm, top_cm=top_cm, width_cm=width_cm, height_cm=height_cm,
                title=title, has_legend=has_legend, has_data_labels=has_data_labels,
            )

        @mcp.tool()
        async def slides_export_pptx(
            path: str,
            output_path: Optional[str] = None,
        ) -> Dict[str, Any]:
            """导出/另存演示文稿为 pptx 文件。output_path 为空时原路径覆盖保存。"""
            return _slides_export_pptx(path, output_path=output_path)
